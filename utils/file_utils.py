from pathlib import Path
from pptx import Presentation

# Ensure data folder exists
Path("data").mkdir(exist_ok=True)

def save_lesson_md(topic: str, lesson_text: str):
    filepath = f"data/{topic.replace(' ', '_')}_lesson.md"
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(lesson_text)
    return filepath

def save_quiz_json(topic: str, quiz_data):
    import json
    filepath = f"data/{topic.replace(' ', '_')}_quiz.json"
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(quiz_data, f, indent=2)
    return filepath

def save_slides_pptx(topic: str, slides_text: str):
    prs = Presentation()
    slides = slides_text.split("**Slide ")[1:]  # split by slides

    for slide_content in slides:
        slide_lines = slide_content.strip().split("\n")
        title_line = slide_lines[0].replace(": Title Slide", "").replace(":", "").strip()
        bullet_points = [line.replace("* ", "") for line in slide_lines if line.startswith("* ")]
        
        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_line if title_line else "Slide"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for bp in bullet_points:
            p = body.add_paragraph()
            p.text = bp

    filepath = f"data/{topic.replace(' ', '_')}_slides.pptx"
    prs.save(filepath)
    return filepath
